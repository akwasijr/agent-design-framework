#!/usr/bin/env python3
"""
Agent Design Framework — Elite Edition
Style: Tech Keynote (Apple/Tesla)
Built following the Elite PowerPoint Designer skill system.

Design System:
  - Colors: Black (#000000), White (#FFFFFF), Accent Blue (#0071E3)
  - Typography: Calibri (titles 72-96pt), Calibri Light (body 24-28pt)
  - Layout: Extreme whitespace, single focal point, 1"+ gutters
  - Max 2 font families, max 4 font sizes
  - One main idea per slide
  - Max 6 lines body text per slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ──────────────────────────────────────────────
# DESIGN SYSTEM — Tech Keynote
# ──────────────────────────────────────────────
BLACK = RGBColor(0x00, 0x00, 0x00)
NEAR_BLACK = RGBColor(0x0A, 0x0A, 0x0A)
DARK_CARD = RGBColor(0x14, 0x14, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x71, 0xE3)       # Apple blue
GRAY_80 = RGBColor(0xCC, 0xCC, 0xCC)
GRAY_60 = RGBColor(0x99, 0x99, 0x99)
GRAY_40 = RGBColor(0x66, 0x66, 0x66)
GRAY_20 = RGBColor(0x33, 0x33, 0x33)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri Light"

# Typography scale (per skill spec)
HERO_SIZE = Pt(84)
SECTION_SIZE = Pt(64)
SLIDE_TITLE = Pt(48)
BODY_LARGE = Pt(32)
BODY = Pt(26)
CAPTION = Pt(18)

# Spacing (1"+ gutters per skill spec)
GUTTER = Inches(1.1)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
CONTENT_W = SLIDE_W - GUTTER * 2  # ~7.8"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────

def black_slide():
    """Add blank slide with black background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def add_text(slide, text, x, y, w, h, font_size=BODY, font_name=BODY_FONT,
             color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a text box with full formatting control."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Set vertical alignment
    tf.paragraphs[0].alignment = align

    p = tf.paragraphs[0]
    p.text = text.split('\n')[0] if '\n' in text else text
    p.font.size = font_size
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing

    # Add remaining lines
    if '\n' in text:
        for line in text.split('\n')[1:]:
            p2 = tf.add_paragraph()
            p2.text = line
            p2.font.size = font_size
            p2.font.name = font_name
            p2.font.color.rgb = color
            p2.font.bold = bold
            p2.font.italic = italic
            p2.alignment = align
            if line_spacing:
                p2.line_spacing = line_spacing

    return txBox


def add_accent_line(slide, x, y, w, thickness=Pt(3)):
    """Add a thin accent-colored line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, w, thickness
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_rect(slide, x, y, w, h, color=DARK_CARD):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def center_x(content_width):
    """Get x position to center content of given width."""
    return (SLIDE_W - content_width) // 2


# ──────────────────────────────────────────────
# SLIDE TYPE TEMPLATES
# ──────────────────────────────────────────────

def title_slide(title, subtitle=None):
    """Hero title slide — 96pt title, centered, accent line."""
    s = black_slide()
    # Top accent line
    add_accent_line(s, Inches(0), Inches(0), SLIDE_W, Pt(4))

    # Title — huge, centered
    add_text(s, title,
             GUTTER, Inches(1.6), CONTENT_W, Inches(2.0),
             font_size=HERO_SIZE, font_name=TITLE_FONT,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    if subtitle:
        add_text(s, subtitle,
                 GUTTER, Inches(3.6), CONTENT_W, Inches(0.8),
                 font_size=BODY_LARGE, font_name=BODY_FONT,
                 color=ACCENT, align=PP_ALIGN.CENTER)

    return s


def chapter_intro(title, subtitle=None):
    """Section divider — large centered text, accent line above."""
    s = black_slide()
    line_w = Inches(1.5)
    add_accent_line(s, center_x(line_w), Inches(1.8), line_w, Pt(3))

    add_text(s, title,
             GUTTER, Inches(2.0), CONTENT_W, Inches(1.8),
             font_size=SECTION_SIZE, font_name=TITLE_FONT,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    if subtitle:
        add_text(s, subtitle,
                 Inches(1.5), Inches(3.8), Inches(7), Inches(0.8),
                 font_size=BODY, font_name=BODY_FONT,
                 color=GRAY_60, align=PP_ALIGN.CENTER)

    return s


def quote_slide(quote_text, attribution=None):
    """Big impactful quote — large text, centered."""
    s = black_slide()
    add_text(s, quote_text,
             GUTTER, Inches(1.2), CONTENT_W, Inches(3.0),
             font_size=Pt(40), font_name=TITLE_FONT,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             line_spacing=Pt(56))

    if attribution:
        add_text(s, attribution,
                 GUTTER, Inches(4.2), CONTENT_W, Inches(0.6),
                 font_size=CAPTION, font_name=BODY_FONT,
                 color=GRAY_60, align=PP_ALIGN.CENTER)

    return s


def key_message_slide(title, points):
    """Title + 2-3 key points with large text."""
    s = black_slide()
    add_text(s, title,
             GUTTER, Inches(0.6), CONTENT_W, Inches(1.0),
             font_size=SLIDE_TITLE, font_name=TITLE_FONT,
             color=WHITE, bold=True, align=PP_ALIGN.LEFT)

    add_accent_line(s, GUTTER, Inches(1.5), Inches(2.0))

    for i, (heading, desc) in enumerate(points):
        y_off = Inches(1.9) + Inches(i * 1.1)
        add_text(s, heading,
                 GUTTER, y_off, CONTENT_W, Inches(0.5),
                 font_size=BODY_LARGE, font_name=TITLE_FONT,
                 color=WHITE, bold=True)
        add_text(s, desc,
                 GUTTER, y_off + Inches(0.45), CONTENT_W, Inches(0.5),
                 font_size=BODY, font_name=BODY_FONT,
                 color=GRAY_60)

    return s


def metrics_dashboard(title, metrics):
    """Big number KPIs — large centered numbers."""
    s = black_slide()
    add_text(s, title,
             GUTTER, Inches(0.5), CONTENT_W, Inches(0.8),
             font_size=SLIDE_TITLE, font_name=TITLE_FONT,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    n = len(metrics)
    card_w = Inches(2.0)
    total_w = card_w * n + Inches(0.3) * (n - 1)
    start_x = center_x(total_w)

    for i, (number, label) in enumerate(metrics):
        x = start_x + (card_w + Inches(0.3)) * i
        # Big number
        add_text(s, str(number),
                 x, Inches(1.8), card_w, Inches(1.5),
                 font_size=Pt(72), font_name=TITLE_FONT,
                 color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        # Label
        add_text(s, label,
                 x, Inches(3.3), card_w, Inches(0.8),
                 font_size=CAPTION, font_name=BODY_FONT,
                 color=GRAY_60, align=PP_ALIGN.CENTER)

    return s


def before_after(title, left_title, left_items, right_title, right_items):
    """Split comparison — muted left, bright right."""
    s = black_slide()
    add_text(s, title,
             GUTTER, Inches(0.5), CONTENT_W, Inches(0.8),
             font_size=SLIDE_TITLE, font_name=TITLE_FONT,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    half_w = Inches(3.6)
    left_x = Inches(0.6)
    right_x = Inches(5.4)

    # Left column (muted)
    add_text(s, left_title,
             left_x, Inches(1.5), half_w, Inches(0.6),
             font_size=BODY_LARGE, font_name=TITLE_FONT,
             color=GRAY_40, bold=True, align=PP_ALIGN.CENTER)
    for i, item in enumerate(left_items):
        add_text(s, item,
                 left_x, Inches(2.2 + i * 0.55), half_w, Inches(0.5),
                 font_size=BODY, font_name=BODY_FONT,
                 color=GRAY_40, align=PP_ALIGN.CENTER)

    # Center divider
    add_rect(s, Inches(4.6), Inches(1.5), Pt(2), Inches(3.5), GRAY_20)

    # Right column (bright)
    add_accent_line(s, right_x, Inches(1.45), half_w, Pt(3))
    add_text(s, right_title,
             right_x, Inches(1.55), half_w, Inches(0.6),
             font_size=BODY_LARGE, font_name=TITLE_FONT,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for i, item in enumerate(right_items):
        add_text(s, item,
                 right_x, Inches(2.2 + i * 0.55), half_w, Inches(0.5),
                 font_size=BODY, font_name=BODY_FONT,
                 color=GRAY_80, align=PP_ALIGN.CENTER)

    return s


def three_column(title, columns):
    """Three column layout — title + 3 blocks."""
    s = black_slide()
    add_text(s, title,
             GUTTER, Inches(0.4), CONTENT_W, Inches(0.8),
             font_size=SLIDE_TITLE, font_name=TITLE_FONT,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    col_w = Inches(2.5)
    gap = Inches(0.3)
    total_w = col_w * 3 + gap * 2
    start_x = center_x(total_w)

    for i, col in enumerate(columns):
        x = start_x + (col_w + gap) * i
        # Accent line on top
        color = col.get('color', ACCENT)
        add_accent_line(s, x, Inches(1.5), col_w, Pt(3))
        # Number or label
        if 'number' in col:
            add_text(s, col['number'],
                     x, Inches(1.7), col_w, Inches(0.9),
                     font_size=Pt(48), font_name=TITLE_FONT,
                     color=color, bold=True, align=PP_ALIGN.CENTER)
        # Title
        add_text(s, col['title'],
                 x, Inches(2.6), col_w, Inches(0.6),
                 font_size=BODY_LARGE, font_name=TITLE_FONT,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Desc
        add_text(s, col['desc'],
                 x + Inches(0.1), Inches(3.2), col_w - Inches(0.2), Inches(1.8),
                 font_size=Pt(20), font_name=BODY_FONT,
                 color=GRAY_60, align=PP_ALIGN.CENTER,
                 line_spacing=Pt(28))

    return s


def bullet_slide(title, items, subtitle=None):
    """Clean bullet hierarchy slide — large bullets."""
    s = black_slide()
    add_text(s, title,
             GUTTER, Inches(0.5), CONTENT_W, Inches(0.8),
             font_size=SLIDE_TITLE, font_name=TITLE_FONT,
             color=WHITE, bold=True)

    if subtitle:
        add_text(s, subtitle,
                 GUTTER, Inches(1.2), CONTENT_W, Inches(0.5),
                 font_size=BODY, font_name=BODY_FONT,
                 color=ACCENT)

    start_y = Inches(1.9) if subtitle else Inches(1.6)
    for i, item in enumerate(items):
        y = start_y + Inches(i * 0.6)
        # Accent dot
        add_rect(s, GUTTER, y + Inches(0.12), Pt(8), Pt(8), ACCENT)
        add_text(s, item,
                 GUTTER + Inches(0.35), y, CONTENT_W - Inches(0.35), Inches(0.5),
                 font_size=BODY, font_name=BODY_FONT,
                 color=GRAY_80)

    return s


def statement_slide(number, principle_title, big_statement, supporting_text):
    """Single design principle as a bold statement."""
    s = black_slide()

    # Number + title row
    add_text(s, number,
             GUTTER, Inches(0.5), Inches(1.0), Inches(0.7),
             font_size=Pt(40), font_name=TITLE_FONT,
             color=ACCENT, bold=True)
    add_text(s, principle_title,
             GUTTER + Inches(1.2), Inches(0.55), Inches(6), Inches(0.6),
             font_size=BODY_LARGE, font_name=BODY_FONT,
             color=GRAY_40)

    # Accent line
    add_accent_line(s, GUTTER, Inches(1.3), Inches(2.0))

    # Big statement
    add_text(s, big_statement,
             GUTTER, Inches(1.7), CONTENT_W, Inches(2.2),
             font_size=Pt(44), font_name=TITLE_FONT,
             color=WHITE, bold=True,
             line_spacing=Pt(58))

    # Supporting text
    add_text(s, supporting_text,
             GUTTER, Inches(4.0), Inches(7), Inches(1.0),
             font_size=BODY, font_name=BODY_FONT,
             color=GRAY_60,
             line_spacing=Pt(36))

    return s


# ──────────────────────────────────────────────
# BUILD THE DECK
# ──────────────────────────────────────────────

# ═══ 1. HERO TITLE ═══
title_slide(
    "Agent Design\nFramework",
    "A Thinking Framework for Agent-Managed Work"
)

# ═══ 2. BIG QUOTE ═══
quote_slide(
    "Every knowledge worker\nbecomes a manager.\n\nThe product they use is\na management tool.",
    "The skill they need is judgment, not execution."
)

# ═══ 3. CHAPTER: THE PROBLEM ═══
chapter_intro("The Problem", "Technology is advancing. Design thinking hasn't kept up.")

# ═══ 4. THREE MISSING LAYERS ═══
key_message_slide("Three Missing Layers", [
    ("Experience", "No design language exists for agent-managed work"),
    ("Accessibility", "Agent creation is still a technical skill"),
    ("Governance", "No standard for boundaries, audits, or responsibility"),
])

# ═══ 5. CHAPTER: THE CORE INSIGHT ═══
chapter_intro("The Core Insight", "From doing to managing")

# ═══ 6. BEFORE/AFTER ═══
before_after(
    "The Paradigm Shift",
    "TODAY",
    ["Process applications", "Create campaigns", "Document encounters", "Grade assignments"],
    "TOMORROW",
    ["Set risk parameters", "Define brand direction", "Oversee care teams", "Design learning journeys"],
)

# ═══ 7. ROLE TRANSFORMATION ═══
bullet_slide("Role Transformation", [
    "Processor  →  Governor",
    "Creator  →  Creative Director",
    "Practitioner  →  Care Architect",
    "Administrator  →  Public Steward",
    "Instructor  →  Learning Architect",
], subtitle="Every role shifts from execution to oversight")

# ═══ 8. CHAPTER: THE EVOLUTION ═══
chapter_intro("The Evolution", "Three phases from code to culture")

# ═══ 9. THREE PHASES ═══
three_column("Three Phases", [
    {"number": "01", "title": "Developer\nTools", "desc": "Agents live in the domain\nof engineers. Code, APIs,\nand frameworks.", "color": ACCENT},
    {"number": "02", "title": "Agent\nPlatforms", "desc": "No-code builders.\nGovernance as product.\nAgent dashboards.", "color": RGBColor(0x8B, 0x5C, 0xF6)},
    {"number": "03", "title": "Agent-Native\nProducts", "desc": "Products designed for\nagent-managed work\nfrom the ground up.", "color": RGBColor(0xF5, 0x9E, 0x0B)},
])

# ═══ 10. CHAPTER: TANGIBLE FUTURES ═══
chapter_intro("Tangible Futures", "What changes — and what emerges")

# ═══ 11. WHAT DISAPPEARS ═══
bullet_slide("What Disappears", [
    "Navigation menus",
    "Search bars",
    "Forms and data entry",
    "Notification overload",
    "Dashboards you check",
    "Settings panels",
])

# ═══ 12. WHAT EMERGES ═══
bullet_slide("What Emerges", [
    "The briefing — curated summary of what matters",
    "Proposal cards — agent recommendations with reasoning",
    "Exception queue — only what agents couldn't handle",
    "Trust dial — adjustable autonomy per agent",
    "Choreography map — agent coordination view",
    "Teaching moments — natural feedback loops",
])

# ═══ 13. KEY METRIC — A DAY WITH AGENTS ═══
metrics_dashboard("A Day With Agents", [
    ("84", "Claims\nProcessed"),
    ("72", "Resolved\nAutonomously"),
    ("5", "Escalated\nto You"),
    ("45m", "Judgment\nWork"),
])

# ═══ 14. BIG STATEMENT ═══
quote_slide(
    "Zero paperwork.",
    "Insurance Claims Manager — Morning to evening"
)

# ═══ 15. THREE INTERACTION MODES ═══
three_column("Three Interaction Modes", [
    {"title": "Delegate", "desc": "Set the goal and\nwalk away. Agent works\nindependently.", "color": ACCENT},
    {"title": "Supervise", "desc": "Agent works, you watch.\nStep in when needed.\nLike managing a junior.", "color": RGBColor(0x8B, 0x5C, 0xF6)},
    {"title": "Collaborate", "desc": "Real-time back-and-forth.\nNeither fully in charge.\nCreative co-creation.", "color": RGBColor(0xF5, 0x9E, 0x0B)},
])

# ═══ 16. CHAPTER: EXPERIENCE DESIGN ═══
chapter_intro("Experience Design", "Interfaces for managing agents, not performing tasks")

# ═══ 17. COMMAND CENTER ═══
key_message_slide("The Command Center", [
    ("One Surface", "Replaces the fragmented multi-app experience"),
    ("Outcomes First", "Show what was achieved, not activity logs"),
    ("Exception Queue", "Judgment, ethics, authority, creativity, empathy"),
])

# ═══ 18. APPROVAL FLOW ═══
key_message_slide("The Approval Flow", [
    ("Agent Proposes", "Recommended action with reasoning and alternatives"),
    ("Human Reviews", "Judge, approve, adjust, or reject"),
    ("Agent Learns", "Every decision refines future behavior"),
])

# ═══ 19. AMBIENT AWARENESS ═══
before_after(
    "Ambient Awareness",
    "HIGH URGENCY",
    ["Low confidence: Interrupt", "High confidence: Notify"],
    "LOW URGENCY",
    ["Low confidence: Queue", "High confidence: Handle silently"],
)

# ═══ 20. STATEMENT — SILENCE ═══
quote_slide(
    "Silence should feel\nlike competence,\nnot absence."
)

# ═══ 21. CHAPTER: DESIGN PRINCIPLES ═══
chapter_intro("Design Principles", "Five rules for agent-managed product design")

# ═══ 22–26. ONE PRINCIPLE PER SLIDE ═══
statement_slide("01", "Outcomes Over Processes",
    "Show what was achieved,\nnot how it was done.",
    "Human attention goes to results.\nProcess details are available on demand, never the default.")

statement_slide("02", "Exceptions Over Routine",
    "Surface what's unusual.\nHide what's normal.",
    "If humans review everything, they might\nas well do the work themselves.")

statement_slide("03", "Trust Building Over Time",
    "Trust is earned,\nnot configured.",
    "Incremental. Evidence-based. Reversible.\nMovement earned through demonstrated competence.")

statement_slide("04", "Values Over Settings",
    "Express boundaries\nas human values.",
    "'Be conservative when unsure.'\nNot risk_score_threshold: 0.65.")

statement_slide("05", "Agent Creation as Document Creation",
    "As natural as making\na spreadsheet.",
    "Domain expertise — not technical skill —\nshould be the requirement.")

# ═══ 27. CHAPTER: INDUSTRIES ═══
chapter_intro("Industries", "Agent-managed work across sectors")

# ═══ 28. INDUSTRY SHIFTS ═══
bullet_slide("Industry Transformation", [
    "Banking — Transaction Processors → Financial Governors",
    "Insurance — Claims Assessors → Risk Adjudicators",
    "Healthcare — Practitioners → Care Orchestrators",
    "Government — Administrators → Citizen Stewards",
    "Education — Instructors → Learning Architects",
    "Manufacturing — Operators → Systems Conductors",
])

# ═══ 29. CHAPTER: EMERGING ═══
chapter_intro("Emerging Considerations", "Questions that will shape the future")

# ═══ 30. EMERGING TOPICS ═══
key_message_slide("What's Coming", [
    ("Agent Identity & Economies", "Credit scores for agents. Machines negotiating with machines."),
    ("Regulation & Emotional Design", "Frameworks for agent decisions. Beyond functional competence."),
    ("Digital Divide 2.0", "Agent access as the next inequality frontier."),
])

# ═══ 31. CLOSING QUOTE ═══
quote_slide(
    "How do humans maintain\nmeaningful control,\nunderstanding, and purpose\nin a world where agents\ndo most of the work?",
    "This is as much philosophical as it is a design question."
)

# ═══ 32. THANK YOU ═══
s = black_slide()
add_accent_line(s, Inches(0), Inches(0), SLIDE_W, Pt(4))
add_accent_line(s, Inches(0), SLIDE_H - Pt(4), SLIDE_W, Pt(4))
add_text(s, "Agent Design\nFramework",
         GUTTER, Inches(1.5), CONTENT_W, Inches(2.0),
         font_size=SECTION_SIZE, font_name=TITLE_FONT,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Designed for the next paradigm of work.",
         GUTTER, Inches(3.6), CONTENT_W, Inches(0.6),
         font_size=BODY, font_name=BODY_FONT,
         color=ACCENT, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# VALIDATION (Step 6 from skill)
# ──────────────────────────────────────────────
print("── Quality Checklist ──")

fonts_used = set()
colors_used = set()
max_text_per_slide = []

for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.font.name:
                    fonts_used.add(para.font.name)
                try:
                    if para.font.color and para.font.color.rgb:
                        colors_used.add(str(para.font.color.rgb))
                except AttributeError:
                    pass
                if para.text.strip():
                    texts.append(para.text.strip())
    max_text_per_slide.append(len(texts))

print(f"  Fonts: {fonts_used} ({'✅' if len(fonts_used) <= 2 else '❌'} max 2)")
print(f"  Colors: {len(colors_used)} used ({'✅' if len(colors_used) <= 8 else '❌'})")
print(f"  Max text items on a slide: {max(max_text_per_slide)}")
print(f"  Total slides: {len(prs.slides)}")

# Check for placeholder text
for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.lower()
                if 'xxxx' in t or 'lorem' in t or 'placeholder' in t:
                    print(f"  ❌ Placeholder found on slide {i}: {para.text[:60]}")

# ──────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────
out = "Agent_Design_Framework_Elite.pptx"
prs.save(out)
print(f"\n✅ Saved: {out} ({len(prs.slides)} slides)")
